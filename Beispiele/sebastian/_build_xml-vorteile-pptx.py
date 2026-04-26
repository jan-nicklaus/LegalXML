"""Generiert xml-vorteile.pptx — PowerPoint-Argumentarium auf Basis
der Kanzlei-Vorlage 'Standard PowerPoint.potx'.

Ausführen:
    uv run --no-project --with python-pptx \\
        python3 Beispiele/sebastian/_build_xml-vorteile-pptx.py
"""
from pathlib import Path
import zipfile, tempfile, os

from pptx import Presentation
from pptx.util import Pt, Cm, Inches, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

REPO = Path(__file__).resolve().parents[2]
TEMPLATE = REPO / "Standard PowerPoint.potx"
OUT = Path(__file__).parent / "xml-vorteile.pptx"

ACCENT = RGBColor(0x38, 0xB0, 0x7D)        # Kanzlei-Grün aus dem Theme
INK    = RGBColor(0x21, 0x21, 0x21)
MUTED  = RGBColor(0x55, 0x55, 0x55)


# --- Vorlage potx → pptx konvertieren (nur content-type ändern) -------
def potx_to_pptx_bytes(potx_path: Path) -> bytes:
    """Konvertiert .potx zu .pptx, indem der Content-Type umgeschrieben
    wird. Inhalt sonst identisch."""
    import io
    buf = io.BytesIO()
    with zipfile.ZipFile(potx_path, 'r') as zin, \
         zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.namelist():
            data = zin.read(item)
            if item == '[Content_Types].xml':
                data = data.replace(
                    b'application/vnd.openxmlformats-officedocument.'
                    b'presentationml.template.main+xml',
                    b'application/vnd.openxmlformats-officedocument.'
                    b'presentationml.presentation.main+xml')
            zout.writestr(item, data)
    return buf.getvalue()


# --- Helper -----------------------------------------------------------
def remove_all_slides(prs):
    """Entfernt alle bestehenden Folien aus der Präsentation. Master-
    Slides und Layouts bleiben erhalten."""
    sldIdLst = prs.slides._sldIdLst
    rId_to_drop = []
    for sldId in list(sldIdLst):
        rId_to_drop.append(sldId.attrib[
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'])
        sldIdLst.remove(sldId)
    for rId in rId_to_drop:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


def find_placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def set_text(placeholder, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    placeholder.text_frame.text = ""
    p = placeholder.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bullets(placeholder, items, *, size=14, level_size_step=2):
    """Füllt einen BODY-Platzhalter mit einer Bullet-Liste.
    items: Liste von (text, level) oder strings."""
    tf = placeholder.text_frame
    tf.text = ""
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        p.text = text
        for run in p.runs:
            run.font.size = Pt(size - level * level_size_step)
            run.font.color.rgb = INK


def add_text_below_title(slide, text, *, top_cm=4, height_cm=12, size=14):
    """Fügt einen Textblock unterhalb des Titels ein (für Layouts ohne
    geeigneten Body-Platzhalter)."""
    left = Cm(2)
    width = Cm(slide.part.package.presentation_part.presentation.slide_width.cm - 4) \
        if False else Cm(21.5)
    box = slide.shapes.add_textbox(left, Cm(top_cm), width, Cm(height_cm))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = INK
    return box


def add_table(slide, data, *, left_cm=1.5, top_cm=4.5, width_cm=22, height_cm=11):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(
        rows, cols, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    ).table

    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(11 if r == 0 else 10)
            run.font.bold = (r == 0 or c == 0)
            if r == 0:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT
            else:
                run.font.color.rgb = INK
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF7, 0xF7, 0xF7)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return tbl


# --- Build ------------------------------------------------------------
def build():
    pptx_bytes = potx_to_pptx_bytes(TEMPLATE)
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    tmp.write(pptx_bytes)
    tmp.close()

    prs = Presentation(tmp.name)
    remove_all_slides(prs)         # Bestand-Folien der Vorlage entfernen
    layouts = list(prs.slide_layouts)

    # Layout-Indices (per Inspektion):
    # 0 Titelfolie | 1 Abschnitts­überschrift | 2 Agenda | 6 Titel & Text
    # 8 Danke | 9 Inhalt mit Überschrift
    L_TITEL    = layouts[0]
    L_ABSCHN   = layouts[1]
    L_AGENDA   = layouts[2]
    L_NUR_TEXT = layouts[4]
    L_INHALT   = layouts[9]
    L_DANKE    = layouts[8]

    # ------------------------------------------------------------------
    # 1 — Titelfolie
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_TITEL)
    # Layout 0 hat zwei BODY-Platzhalter (idx 10 und 14).
    # idx 10 = Titelzeile (oberer Bereich), idx 14 = Subtitle.
    ph10 = find_placeholder(s, 10)
    ph14 = find_placeholder(s, 14)
    if ph10:
        set_text(ph10, "XML statt PDF/DOCX", size=40, bold=True, color=ACCENT)
    if ph14:
        set_text(ph14, "Strukturierte Rechtsdokumente — ein Argumentarium",
                 size=22, color=INK)

    # ------------------------------------------------------------------
    # 2 — Agenda
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_AGENDA)
    ph14 = find_placeholder(s, 14)
    if ph14:
        add_bullets(ph14, [
            "Sinn und Zweck",
            "Drei Formate auf einen Blick",
            "Warum PDF und DOCX nicht reichen",
            "Vorteile von Domänen-XML",
            "Was XML nicht ist",
            "Faustregel",
        ], size=20)

    # ------------------------------------------------------------------
    # 3 — Sinn und Zweck (Abschnittsüberschrift)
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_ABSCHN)
    ph0 = find_placeholder(s, 0)
    if ph0:
        set_text(ph0, "Sinn und Zweck", size=44, bold=True, color=ACCENT)

    # ------------------------------------------------------------------
    # 4 — Sinn und Zweck (Inhalt)
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_INHALT)
    title = find_placeholder(s, 0)
    body  = find_placeholder(s, 2) or find_placeholder(s, 1)
    if title:
        set_text(title, "Wozu strukturierte Rechtsdokumente?", size=28, bold=True, color=ACCENT)
    if body:
        add_bullets(body, [
            "Rechtsschriften und Urteile bestehen mehrheitlich aus Text.",
            "Im Verfahren wird wechselseitig Bezug genommen — auf inhaltlich verwandte Stellen.",
            ("Sachverhaltselemente korrespondieren mit Sachverhaltselementen,", 1),
            ("rechtliche Argumente werden gegen rechtliche Argumente erwidert.", 1),
            "Strukturierte Daten machen diese Bezüge maschinell nachvollziehbar.",
            "Strukturierung bewusst minimal — gerade so präzise wie nötig.",
            "Externe Ressourcen (Beweismittel, Akten) ebenfalls strukturiert verweisbar.",
        ], size=16)

    # ------------------------------------------------------------------
    # 5 — Drei Formate auf einen Blick
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_INHALT)
    title = find_placeholder(s, 0)
    if title:
        set_text(title, "Drei Formate auf einen Blick", size=28, bold=True, color=ACCENT)

    table_data = [
        ["", "PDF", "DOCX", "Domänen-XML"],
        ["Modelliert",            "Pixel / Glyphen",     "Textverarbeitung (XML)",  "Juristische Inhalte"],
        ["Schema",                "—",                   "OOXML (Layout)",          "Eigen (Erwägung, …)"],
        ["Validierbar",           "—",                   "Layout-Struktur",         "Pflichtfelder, Formate"],
        ["Maschinell auswertbar", "Nur über OCR",        "Über Stilvorlagen",       "Direkt, eindeutig"],
        ["Versionssichtbar",      "Nein",                "Binär (ZIP)",             "Zeile für Zeile"],
        ["Langzeitlesbar",        "Versionsabhängig",    "OOXML offen, breit",      "Offen, schmal"],
    ]
    add_table(s, table_data, top_cm=4.5, height_cm=11)

    # ------------------------------------------------------------------
    # 6 — Warum PDF und DOCX nicht reichen (Abschnittsüberschrift)
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_ABSCHN)
    ph0 = find_placeholder(s, 0)
    if ph0:
        set_text(ph0, "Warum PDF und DOCX nicht reichen", size=36, bold=True, color=ACCENT)

    # ------------------------------------------------------------------
    # 7 — PDF: Layout statt Inhalt
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_INHALT)
    title = find_placeholder(s, 0)
    body  = find_placeholder(s, 2) or find_placeholder(s, 1)
    if title:
        set_text(title, "PDF — Layout statt Inhalt", size=28, bold=True, color=ACCENT)
    if body:
        add_bullets(body, [
            "Sandwich-PDFs: Bild + OCR-Textschicht — Diskrepanzen unbemerkt.",
            "Glyphen statt Bedeutung: «Erwägung Nr. 3» nur visuell erkennbar.",
            "Spezifikations-Unschärfe: dieselbe Information vielfältig darstellbar.",
            "Inhaltliche Validierung nicht vorgesehen.",
            "Hoher Metadaten-Anteil (Schriftarten, Sicherheit, Rendering).",
        ], size=18)

    # ------------------------------------------------------------------
    # 8 — DOCX: XML, aber das falsche
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_INHALT)
    title = find_placeholder(s, 0)
    body  = find_placeholder(s, 2) or find_placeholder(s, 1)
    if title:
        set_text(title, "DOCX — XML, aber das falsche", size=28, bold=True, color=ACCENT)
    if body:
        add_bullets(body, [
            "DOCX ist intern XML (Office Open XML, ZIP-Container).",
            "Aber: generisches Textverarbeitungs-Modell — kein Domänen-Modell.",
            ("Konzepte: Absatz, Stilvorlage, Tabelle.", 1),
            ("Nicht: Erwägung, Geschäftsnummer, Rechtsbegehren.", 1),
            "OOXML validiert die Layout-Struktur, nicht den juristischen Inhalt.",
            "ZIP-Container — nicht direkt diff-bar in Versionskontrolle.",
        ], size=18)

    # ------------------------------------------------------------------
    # 9 — Vorteile (Abschnittsüberschrift)
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_ABSCHN)
    ph0 = find_placeholder(s, 0)
    if ph0:
        set_text(ph0, "Vorteile von Domänen-XML", size=40, bold=True, color=ACCENT)

    # ------------------------------------------------------------------
    # 10 — Operative Vorteile
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_INHALT)
    title = find_placeholder(s, 0)
    body  = find_placeholder(s, 2) or find_placeholder(s, 1)
    if title:
        set_text(title, "Sieben operative Vorteile", size=28, bold=True, color=ACCENT)
    if body:
        add_bullets(body, [
            "Maschinell auswertbar ohne Heuristik (XPath statt OCR-Regex).",
            "Verlustfreie Beziehungen (Vorinstanz, Weiterzug, Zitate).",
            "Validierung vor dem Versand (Pflichtfelder, Geschäftsnummer-Format).",
            "Mehrere Darstellungen aus einer Quelle (PDF, HTML, JSON, TXT).",
            "Versionierung mit Diff-Sinn — inhaltliche Änderungen sichtbar.",
            "Anonymisierung als Datenfeld (steuerbar, maschinell prüfbar).",
            "Stabile Anker für Querverweise (überleben Textänderungen).",
        ], size=16)

    # ------------------------------------------------------------------
    # 11 — Strategischer Mehrwert
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_INHALT)
    title = find_placeholder(s, 0)
    body  = find_placeholder(s, 2) or find_placeholder(s, 1)
    if title:
        set_text(title, "Strategischer Mehrwert", size=28, bold=True, color=ACCENT)
    if body:
        add_bullets(body, [
            "Statistische Auswertbarkeit über Korpora.",
            "KI-Tauglichkeit (hochwertige Trainingsgrundlage).",
            "Anschluss an eJustiz-Standards (Akoma Ntoso, Justitia 4.0).",
            "Langzeitarchivierung — offen, herstellerunabhängig.",
            "Kanzleiweite Konsistenz durch erzwungene Pflichtfelder.",
        ], size=18)

    # ------------------------------------------------------------------
    # 12 — Was XML nicht ist
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_INHALT)
    title = find_placeholder(s, 0)
    body  = find_placeholder(s, 2) or find_placeholder(s, 1)
    if title:
        set_text(title, "Was XML nicht ist", size=28, bold=True, color=ACCENT)
    if body:
        add_bullets(body, [
            "Kein Lesedokument — braucht Viewer oder generierte Ansicht.",
            "Kein Schreibwerkzeug — Editor erzeugt XML im Hintergrund.",
            "Kein Layout-Format — PDF wird weiter gebraucht, aber generiert aus XML.",
            "Kein Allheilmittel für bildhafte Beilagen (Pläne, Fotos, Gutachten).",
        ], size=18)

    # ------------------------------------------------------------------
    # 13 — Faustregel
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(L_NUR_TEXT)
    ph14 = find_placeholder(s, 14)
    if ph14:
        tf = ph14.text_frame
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "XML als Quellformat,"
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = ACCENT

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = "PDF und DOCX als Ausgabeformate."
        r2.font.size = Pt(36)
        r2.font.bold = True
        r2.font.color.rgb = ACCENT

        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = ""
        r3.font.size = Pt(8)

        p4 = tf.add_paragraph()
        p4.alignment = PP_ALIGN.CENTER
        r4 = p4.add_run()
        r4.text = ("Strukturierte Daten für maschinelle Anwendungsfälle, "
                   "gewohntes Layout für die Akte.")
        r4.font.size = Pt(18)
        r4.font.color.rgb = MUTED

    # ------------------------------------------------------------------
    # 14 — Danke
    # ------------------------------------------------------------------
    prs.slides.add_slide(L_DANKE)

    # Speichern
    prs.save(str(OUT))
    os.unlink(tmp.name)
    print(f"Geschrieben: {OUT}  ({OUT.stat().st_size:,} Bytes, "
          f"{len(prs.slides)} Folien)")


if __name__ == "__main__":
    build()
